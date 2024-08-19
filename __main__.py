#! /usr/bin/env python

# Standard library imports
import logging
from pathlib import Path
import re

# Third-party library imports
from matplotlib.colors import ListedColormap
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import seaborn as sns
import yaml

# Configure seaborn
sns.set_style(style="whitegrid")
sns.set_palette("muted")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger()

# Define the regular expression pattern to match illegal characters in Windows file paths
ILLEGAL_WINDOWS_PATH_CHARACTERS = re.compile(r'[\\/:*?"<>|\n]')

# Define the regular expression pattern to match the APG No without the pattern
APG_NO_PATTERN = re.compile(r'(\w+\.\d+)')
ENCODING = 'utf-8'

# Define the directory paths
CODING_DIRECTORY = Path(__file__).parent
MAIN_DIRECTORY = CODING_DIRECTORY.parent
GRAPHICS_DIRECTORY = MAIN_DIRECTORY / "Grafikler"
REPORTS_DIRECTORY = MAIN_DIRECTORY / "Raporlar"

# Define the file paths
CONFIG_PATH = CODING_DIRECTORY / "config.yaml"
PRESENTATION_INTRO_TEMPLATE_PATH = CODING_DIRECTORY / "presentation_intro_template.txt"

# Load YAML configuration
with CONFIG_PATH.open(encoding=ENCODING) as config_file:
    config = yaml.safe_load(config_file)

MASTER_FILE = MAIN_DIRECTORY / config['MASTER_FILE']
TEMPLATE_PATH = MAIN_DIRECTORY / config['TEMPLATE_PATH']

COMPANY_GROUPS = config['COMPANY_GROUPS']
COMPANY_GROUPS_EXCLUDED_FROM_REPORT = config['COMPANY_GROUPS_EXCLUDED_FROM_REPORT']

# Check if the companies in COMPANY_GROUPS_EXCLUDED_FROM_REPORT are also in COMPANY_GROUPS
if set(COMPANY_GROUPS_EXCLUDED_FROM_REPORT) - set(COMPANY_GROUPS):
    raise ValueError("COMPANY_GROUPS_EXCLUDED_FROM_REPORT contains companies that are not in COMPANY_GROUPS. Please check the config.yaml file.")

NUM_OF_COMPANIES = sum(len(companies) for companies in COMPANY_GROUPS.values())
COMPANIES_RANGE = np.arange(1, NUM_OF_COMPANIES + 1)

REPORT_TYPE_CHOICES = "yariyillik", "yillik"
REPORT_YEAR = config['REPORT_YEAR']
REPORT_TYPE = config['REPORT_TYPE']

SIGMA: int = config['SIGMA']
START_COL = 3
END_COL = START_COL + NUM_OF_COMPANIES
DEFAULT_DECIMAL_DIGITS = config['DEFAULT_DECIMAL_DIGITS']

# Define the company color indicators
GROUP_COMPANY_INDICATOR = 0
RIVAL_COMPANY_INDICATOR = 1

# PLOT PARAMETERS
ANNOTATION_OFFSET_PIXELS = tuple(config['ANNOTATION_OFFSET_PIXELS'])
ANNOTATION_FONT_SIZE = config['ANNOTATION_FONT_SIZE']
HORIZONTAL_MEAN_COLOR = config['HORIZONTAL_MEAN_COLOR']
HORIZONTAL_MEAN_ALPHA = config['HORIZONTAL_MEAN_ALPHA']
FONT_SIZE = config['FONT_SIZE']
OVERLAY_GRAPH_COLOR_MAP = ListedColormap(config['OVERLAY_GRAPH_COLOR_MAP'])
OVERLAY_GRAPH_BAR_COLOR = config['OVERLAY_GRAPH_BAR_COLOR']


def generate_presentation_intro_text(company_list: list[str]) -> str:
    """
    Generate the introductory text for the presentation based on the company group.
    company_list refers to the list of companies in the current company group.
    This will be used in the first slide of the presentation.
    """
    # Load the presentation intro template from the template text file
    presentation_text_template = PRESENTATION_INTRO_TEMPLATE_PATH.read_text(encoding=ENCODING)

    # Iterate through the company list and generate the text for each company in the company group
    company_lines = []
    for i, company in enumerate(company_list, start=1):
        line = f"- {i} numaralı Şirket, {company}'ı"
        # Add the appropriate suffix for the last company in the list
        if i == len(company_list):
            line += " temsil etmekte iken"
        company_lines.append(line)

    # Add the indicator for the remaining companies that doesn't belong to the group
    company_lines.append(f"{num_of_group_companies + 1} - {NUM_OF_COMPANIES}")
    company_enumeration_text = "\n".join(company_lines)

    # Replace the placeholders in the template with the generated text
    return presentation_text_template.format(
        company_text=company_enumeration_text,
        company_group=company_group,
        num_of_APG=unique_categories_amount
    )


def filtered_mean(row: pd.Series) -> int:
    """
    Calculate the mean of each row within a specified number of standard deviations
    and add it to a new column named 'filtered_mean'.

    Parameters:
    df (pandas.DataFrame): The DataFrame containing the data.
    start_col (int): The starting column index for the range of columns to consider.
    end_col (int): The ending column index (exclusive) for the range of columns to consider.
    sigma (int, optional): The number of standard deviations to include. Default is 2.

    Returns:
    pandas.DataFrame: The DataFrame with the new column 'filtered_mean'.
    """
    # Extract the relevant columns
    data = row[START_COL:END_COL]

    # Compute the mean and standard deviation
    mean = data.mean()
    std = data.std()

    # Define the range within the specified number of standard deviations
    lower_bound = mean - SIGMA * std
    upper_bound = mean + SIGMA * std

    # Filter the data to include only values within the range
    filtered_data = data[(data >= lower_bound) & (data <= upper_bound)]

    # Compute the mean of the filtered values
    return filtered_data.mean()



def format_percentage(value: float, decimal_digits: int = DEFAULT_DECIMAL_DIGITS) -> str:
    """
    Format a float value as a percentage with 'DEFAULT_DECIMAL_DIGITS' decimal points and the percent sign in front.

    Parameters:
    value (float): The float value to be formatted as a percentage.
    decimal_digits (int): The number of decimal digits to include. Default is 'DEFAULT_DECIMAL_DIGITS'.

    Returns:
    str: The formatted percentage string.
    """
    return "%{:.{}f}".format(value * 100, decimal_digits).replace('.', ',')


def shuffle_columns(df: pd.DataFrame, company_list: list[str]) -> pd.DataFrame:
    """
    Shuffle the columns of a DataFrame, excluding the first three and the last one.
    Place the company_list columns at the beginning of the shuffle range.
    """
    # Fixed columns are the first three and the last one
    fixed_columns = df.columns[:START_COL].tolist() + df.columns[END_COL:].tolist()

    # Columns to be shuffled, excluding the company_list columns
    columns_to_shuffle = [col for col in df.columns[START_COL:END_COL] if col not in company_list]

    # Shuffle the columns
    shuffled_columns = np.random.permutation(columns_to_shuffle)

    # New column order: first the fixed columns, then the company_list, then the shuffled columns
    new_column_order = fixed_columns[:START_COL] + company_list + shuffled_columns.tolist() + fixed_columns[START_COL:]

    group_df = df[new_column_order]
    group_df.columns.values[START_COL:END_COL] = COMPANIES_RANGE

    return group_df


def standardgraph(row: pd.Series) -> plt.Figure:
    """
    Create a standard scatter plot for the given row.
    """
    # If next line doesn't exist, the graph will be overwritten by each row. Somehow necessary.
    ax = plt.figure()

    # Create the scatter plot
    ax = sns.scatterplot(
        data=transposed,
        x="companies",
        y=row['APG No'],
        legend=False,
        hue=company_color_indicator
    )
    ax.set(xlabel=None)
    ax.set(ylabel=None)
    ax.grid(axis='y')
    ax.set_xticks(COMPANIES_RANGE)
    ax.set_xticklabels(COMPANIES_RANGE)
    ax.set(title=row['APG Full Name'])

    # annotate points on the graph
    for company_index in COMPANIES_RANGE:
        value = row[company_index]
        formatted_value = format_percentage(value) if row["Birim"] == "%" else str(value)
        plt.annotate(
            text=formatted_value,
            xy=(company_index, value),
            fontsize=ANNOTATION_FONT_SIZE,
            xytext=ANNOTATION_OFFSET_PIXELS,
            textcoords="offset pixels"
        )

    # set the y-label to indicate that the corresponding KPI is in TL
    if row["Birim"] == "TL":
        ax.set(ylabel='TL')

    # Set y-axis tick labels if needed
    if row["Birim"] == "%":
        ax.set_yticks(ax.get_yticks())
        ax.set_yticklabels(map(format_percentage, ax.get_yticks()))

    # draw horizontal mean value
    ax.axhline(
        y=row["filtered_mean"],
        color=HORIZONTAL_MEAN_COLOR,
        alpha=HORIZONTAL_MEAN_ALPHA,
    )

    return ax.get_figure()

def stackedgraph(row: pd.Series) -> plt.Figure:
    """
    Create a stacked bar plot for the given row. The bars are stacked based on the Category No.
    Values are formatted as percentages and the total percentage will add up to 100%.
    """
    stacked_apg_nos = category_to_apg_dict[row["Category No"]]
    legend_labels = category_to_apg_full_name_dict[row["Category No"]]

    # Create a figure and axis
    fig, ax = plt.subplots()

    # Plot the transposed data for the specific categories, with stacked bars
    transposed[stacked_apg_nos].plot(kind='bar', stacked=True, ax=ax)

    # Update the legend to be at the bottom of the graph
    ax.legend(labels=legend_labels, loc='center left', bbox_to_anchor=(1, 0.5))

    # Customize the x-ticks and their labels
    ax.set_xticks(COMPANIES_RANGE)
    ax.set_xticklabels(COMPANIES_RANGE)

    # Customize the y-ticks and their labels, formatting them as percentages
    ax.set_yticks(ax.get_yticks())
    ax.set_yticklabels(map(format_percentage, ax.get_yticks()))

    return fig


def overlayedgraph(row: pd.Series) -> plt.Figure:
    """
    Create an overlayed bar and scatter plot for the given row
    """
    apg = row["APG Group"]
    alt_bilgi = f"{apg} EK"
    ara_df = shuffled_df[shuffled_df["APG Group"] == apg]
    main_row = ara_df[ara_df["APG No"] == apg].iloc[0]
    ek_row = ara_df[ara_df["APG No"] == alt_bilgi].iloc[0]

    # bar plot (sub-info)
    fig, ax1 = plt.subplots(figsize=(8, 4))
    ax1.set_xticks(COMPANIES_RANGE)
    ax1.set_xticklabels(COMPANIES_RANGE)
    ax1.bar(
        x=transposed.companies,
        height=transposed[alt_bilgi],
        color=OVERLAY_GRAPH_BAR_COLOR,
        width=0.4,
        zorder=2
    )
    ax1.grid(visible=False)
    ax1.set(ylabel=f"{ek_row['APG İsmi']}\n(Çubuk Gösterim)")
    ax1.set_yticks(ax1.get_yticks())
    ax1.set_yticklabels(map(format_percentage if ek_row["Birim"] == "%" else str, ax1.get_yticks()))

    for x, y in zip(transposed.companies, transposed[alt_bilgi]):
        formatted_text = format_percentage(y) if ek_row["Birim"] == "%" else str(y)
        plt.text(
            x,
            0,
            formatted_text,
            horizontalalignment='center',
            verticalalignment='bottom'
        )

    # create the second (scatter) graph
    ax2 = ax1.twinx()
    ax2.scatter(
        x=transposed.companies,
        y=transposed[main_row['APG No']],
        c=company_color_indicator,
        cmap=OVERLAY_GRAPH_COLOR_MAP,
        zorder=3
    )

    for company_index in COMPANIES_RANGE:
        value = main_row[company_index]
        formatted_value = format_percentage(value) if main_row["Birim"] == "%" else str(value)
        plt.annotate(
            text=formatted_value,
            xy=(company_index, value),
            xytext=(2,5),
            textcoords="offset pixels",
            fontsize=ANNOTATION_FONT_SIZE,
            fontweight='normal',
            zorder=4
        )
    ax2.tick_params(axis='y', labelsize=ANNOTATION_FONT_SIZE)
    ax2.grid(axis='y')
    ax2.set(ylabel=f"{main_row['APG İsmi']}\n(Nokta Gösterim)")
    ax2.set_yticks(ax2.get_yticks())
    ax2.set_yticklabels(map(format_percentage if main_row["Birim"] == "%" else str, ax2.get_yticks()))
    # draw horizontal mean value
    ax2.axhline(
        y=main_row["filtered_mean"],
        color=HORIZONTAL_MEAN_COLOR,
        alpha=HORIZONTAL_MEAN_ALPHA,
    )

    return fig

def create_powerpoint():
    """
    Create a PowerPoint presentation with the graphs for the given company group.
    """
    graphics_save_directory = GRAPHICS_DIRECTORY / f'{REPORT_YEAR}_{REPORT_TYPE}' / company_group
    # Create the directories if they don't exist
    graphics_save_directory.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(TEMPLATE_PATH)

    bulgu_shapes = [
        shape for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text.startswith("Bulgu")
    ]
    if len(bulgu_shapes) != shuffled_df['Bulgu?'].sum():
        raise ValueError("The number of 'Bulgu' shapes in the presentation does not match the number of 'Bulgu' rows in the DataFrame.")

    bulgu_iterator = iter(bulgu_shapes)

    for _, row in shuffled_df.iterrows():
        # Call the appropriate function based on grafik_tipi
        grafik_tipi = row["Grafik_tipi"]
        if grafik_tipi == "standard":
            fig = standardgraph(row)
        elif grafik_tipi == "stacked":
            fig = stackedgraph(row)
        elif grafik_tipi == "overlayed":
            fig = overlayedgraph(row)
        else:
            raise ValueError(f"Unknown grafik_tipi: {grafik_tipi}")

        # save figure to local directory
        apg_pic_name = f'{row["APG Full Name"]}.png'
        # Replace illegal characters in the file name with underscores
        clean_apg_path = ILLEGAL_WINDOWS_PATH_CHARACTERS.sub('_', apg_pic_name)
        pic_path = graphics_save_directory / clean_apg_path
        # Create the directories if they don't exist
        pic_path.parent.mkdir(parents=True, exist_ok=True)
        # Save the figure to the local directory
        fig.savefig(pic_path, bbox_inches='tight')
        logger.info(f"Saved the figure for {row['APG No']} to {pic_path.relative_to(MAIN_DIRECTORY)}")
        # Close the figure to free up resources
        plt.close(fig)

        # Add the figure to the slide
        slide = presentation.slides[row["Sayfa"] - 1]  # Because of 0-based indexing
        left, top, height, width = row["Left"], row["Top"], row["Height"], row["Width"]
        slide.shapes.add_picture(str(pic_path), left, top, width, height)

        # Add the bulgu text to the slide
        if row["Bulgu?"]:
            shape = next(bulgu_iterator)
            filtered_mean_value = row["filtered_mean"]
            formatted_filtered_mean = format_percentage(filtered_mean_value) if row["Birim"] == "%" else str(filtered_mean_value)
            shape.text = f'Bulgu\n{NUM_OF_COMPANIES} şirketin ortalaması {formatted_filtered_mean} olarak tespit edilmiştir.'
            paragraphs = shape.text_frame.paragraphs
            paragraphs[0].font.bold = True
            for paragraph in paragraphs:
                paragraph.font.size = Pt(FONT_SIZE)

    # ilk ve ikinci slaytda yılı ve şirket ismini değiştirme
    ay = 6 if REPORT_TYPE == REPORT_TYPE_CHOICES[0] else 12
    presentation.slides[0].shapes[3].text = f'{REPORT_YEAR} Yılı {ay} Aylık Döneme Ait Performans Göstergesi Sonuçları'
    presentation.slides[1].shapes[3].text = generate_presentation_intro_text(company_list)
    presentation_path = REPORTS_DIRECTORY / f'{REPORT_YEAR}_{REPORT_TYPE}' / company_group / f'Kıyaslama Raporu {REPORT_YEAR}_{REPORT_TYPE}_{company_group}.pptx'
    # Create the directories if they don't exist
    presentation_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(presentation_path)

    logger.info(f"Saved the presentation for {company_group} to {presentation_path.relative_to(MAIN_DIRECTORY)}")

# This line ensures that the code is only run when the script is executed directly, not when it is imported as a module.
if __name__ == "__main__":
    # Read the data from the Excel file and divide it into two DataFrames for the data and the layout
    dataframe_dict = pd.read_excel(MASTER_FILE, sheet_name=[f"{REPORT_YEAR}_Total_Veriler", "pptx_layout"])
    df = dataframe_dict[f"{REPORT_YEAR}_Total_Veriler"]
    pptx_layout = dataframe_dict["pptx_layout"]

    # Merge two sheets on their APG No
    merged_df = pd.merge(df, pptx_layout, left_on='APG No', right_on='APG Kodu', how='left')

    # Make sure that the Sayfa column is an integer, otherwise it complains because automatically it assigns float
    merged_df.Sayfa = merged_df.Sayfa.astype(int)

    # Add Category from APG No
    merged_df['Category No'] = merged_df['APG No'].str.split('.').str[0]
    unique_categories_amount = merged_df['Category No'].nunique()
    # Add a new column to the DataFrame that concatenates the APG No and APG İsmi
    merged_df['APG Full Name'] = merged_df.apply(lambda row: f'{row["APG No"]}-{row["APG İsmi"]}', axis=1)
    # Extract subcategory using regex to group EK APG No's into a common category of APG No's
    merged_df['APG Group'] = merged_df['APG No'].str.extract(APG_NO_PATTERN)[0]

    # For the stacked graph, create a dictionary that maps each category to a list of APG No's
    stacked_df = merged_df[merged_df.Grafik_tipi == "stacked"][["Category No", "APG No", "APG Full Name"]]
    category_to_apg_dict = stacked_df.groupby('Category No')['APG No'].apply(list).to_dict()
    category_to_apg_full_name_dict = stacked_df.groupby('Category No')['APG Full Name'].apply(list).to_dict()

    # Create reports for each company group
    for company_group, company_list in COMPANY_GROUPS.items():
        if company_group in COMPANY_GROUPS_EXCLUDED_FROM_REPORT:
            continue
        # Create a list to indicate a company group (0) or their rivals (1) for each company group
        num_of_group_companies = len(company_list)
        num_of_other_companies = NUM_OF_COMPANIES - num_of_group_companies
        company_color_indicator = [GROUP_COMPANY_INDICATOR] * num_of_group_companies + [RIVAL_COMPANY_INDICATOR] * num_of_other_companies

        # Shuffle columns for each group and store in a list
        shuffled_groups = [
            shuffle_columns(group, company_list)
            for _, group
            in merged_df.groupby('Category No', sort=False)  # Unless the sort is False, APGs are grouped in lexicographical order, not in the order they appear in the Excel file
        ]

        # Merge the shuffled groups back together and reset the column names
        shuffled_df = pd.concat(shuffled_groups).reset_index(drop=True)

        # Apply the filtered_mean function to each row and create a new column
        shuffled_df['filtered_mean'] = shuffled_df.apply(filtered_mean, axis=1)
        # shuffled_df.to_excel(REPORTS_DIRECTORY / f'{REPORT_YEAR}_{REPORT_TYPE}' / company_group / f"{REPORT_YEAR}_{REPORT_TYPE}_{company_group}_Shuffled.xlsx", index=False)

        # Select the relevant columns and transpose the DataFrame and reset the index to companies
        # Will be used for the graphs
        transposable = shuffled_df[["APG No"] + list(COMPANIES_RANGE)]
        transposed = transposable.set_index("APG No").T.reset_index().rename(columns={"index": "companies"})

        create_powerpoint()
